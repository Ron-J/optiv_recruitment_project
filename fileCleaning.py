import os
import shutil
from PIL import Image
import pytesseract
from presidio_image_redactor import ImageRedactorEngine
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import fitz  # PyMuPDF for PDF
from pptx import Presentation
import openpyxl

# Configure Tesseract path (adjust if needed)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Initialize Presidio engines
analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()


def clean_image(file_path: str):
    """Redact PII from an image."""
    image = Image.open(file_path)
    redactor = ImageRedactorEngine()
    redacted_image = redactor.redact(image=image)
    redacted_image.save(file_path)  # overwrite original


import fitz
import os
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine

analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

def clean_pdf(file_path: str):
    """Redact PII from a PDF preserving layout as much as possible."""
    doc = fitz.open(file_path)

    for page in doc:
        # A list to store redaction annotations
        redact_list = []
        blocks = page.get_text("blocks")

        for block in blocks:
            x0, y0, x1, y1, text, *_ = block

            if not text.strip():
                continue

            # Assuming analyzer and anonymizer are defined globally or passed
            # ...
            # Get PII analysis results
            results = analyzer.analyze(text=text, language="en")
            anonymized = anonymizer.anonymize(text=text, analyzer_results=results).text
            # ...

            if text != anonymized:
                rect = fitz.Rect(x0, y0, x1, y1)
                
                # 1. Add a redaction annotation
                # Use a block color (e.g., black (0, 0, 0) or white (1, 1, 1)) and the replacement text
                page.add_redact_annot(
                    rect, 
                    text=anonymized, 
                    fill=(1, 1, 1), # White fill color for the redacted area
                    text_color=(0, 0, 0), # Black text color for the replacement text
                    fontsize=11
                )
                
        # 2. Apply all redactions on the page
        # This is the crucial step that removes the original content and applies the replacement
        page.apply_redactions() 

    cleaned_path = file_path.replace(".pdf", "_cleaned.pdf")
    doc.save(cleaned_path, deflate=True)
    doc.close()

    # Replace original file safely
    import os # Ensure os is imported
    os.replace(cleaned_path, file_path)



def clean_ppt(file_path: str):
    """Redact PII from a PowerPoint file, including tables."""
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            # Case 1: Normal text boxes / placeholders
            if shape.has_text_frame:
                text = shape.text
                if text.strip():
                    results = analyzer.analyze(text=text, language="en")
                    anonymized = anonymizer.anonymize(text=text, analyzer_results=results).text
                    shape.text = anonymized

            # Case 2: Tables
            if shape.shape_type == 19:  # 19 == MSO_SHAPE_TYPE.TABLE
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text = cell.text
                            results = analyzer.analyze(text=text, language="en")
                            anonymized = anonymizer.anonymize(text=text, analyzer_results=results).text
                            cell.text = anonymized

    prs.save(file_path)


def clean_excel(file_path: str):
    """Redact PII from an Excel file."""
    wb = openpyxl.load_workbook(file_path)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    results = analyzer.analyze(text=cell.value, language="en")
                    anonymized = anonymizer.anonymize(text=cell.value, analyzer_results=results).text
                    cell.value = anonymized
    wb.save(file_path)


def clean_file(file_path: str):
    """Detect file type and clean accordingly."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext in [".png", ".jpg", ".jpeg"]:
        clean_image(file_path)
    elif ext == ".pdf":
        clean_pdf(file_path)
    elif ext == ".pptx":
        clean_ppt(file_path)
    elif ext in [".xls", ".xlsx"]:
        clean_excel(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# Example usage:
if __name__ == "__main__":
    test_files = [
        "test.jpeg",
        "sample.pdf",
        "slides.pptx",
        "data.xlsx"
    ]

    for f in test_files:
        if os.path.exists(f):
            print(f"Cleaning {f}...")
            clean_file(f)
            print(f"{f} cleaned successfully ✅")
