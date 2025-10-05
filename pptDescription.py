from pptx import Presentation
from GeminiAPI import print_file_description_and_key_findings
from fileCleaning import clean_ppt
def process_PPT(file_path):
    clean_ppt(file_path)
    prs = Presentation(file_path)
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_list = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                text_list.append(shape.text.strip())
        slide_text = "\n".join(text_list).strip()
        slides_text.append({"slide_index": idx, "text": slide_text})

    prompt=f"""
    Give a title and the file desciption for the PPT file with the following content in about 30 words:
    {slides_text}"""

    return print_file_description_and_key_findings(prompt)